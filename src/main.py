# src/main.py

import os
import subprocess
import sys
import logging
from importlib import metadata

from preprocessing import preprocess_text, text_statistics, count_syllables_pt
from analysis import (
    extract_entities,
    extract_pos_tags,
    dependency_parsing,
    keyword_extraction,
    extract_relationships,
    lda_topic_modeling,
    sentiment_analysis,
    analyze_connectors,
    spelling_correction,
    readability_scores,
    extract_dates,
    extract_actions_and_responsibles,
    check_verb_agreement,
    detect_person_changes
)
from visualization import (
    generate_word_cloud,
    generate_entity_network,
    generate_dense_pixel_display,
    generate_topic_visualization,
    generate_action_flow
)
from database import store_data_in_database
from gui import TextMiningGUI
from utils import download_nltk_packages, load_spacy_model

import tkinter as tk
from tkinter import filedialog, messagebox

from datetime import datetime
import logging

def setup_logging(output_folder):
    """
    Configura o logging para o projeto.
    """
    try:
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
            print(f"Pasta '{output_folder}' criada para armazenar logs.")
            logging.info(f"Pasta '{output_folder}' criada para armazenar logs.")
        
        log_filename = f"execucao_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
        log_filepath = os.path.join(output_folder, log_filename)
        logging.basicConfig(
            filename=log_filepath,
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
        )
        logging.info("Configuração de logging inicializada.")
        print("Configuração de logging inicializada.")
    except Exception as e:
        print(f"Erro ao configurar logging: {e}")
        sys.exit(1)

def read_document(file_path):
    """
    Lê o conteúdo de um documento com base na sua extensão.
    
    Parâmetros:
        file_path (str): Caminho para o arquivo.
    
    Retorna:
        str: Texto extraído do documento ou uma string vazia em caso de erro.
    """
    try:
        if not os.path.isfile(file_path):
            logging.error(f"Arquivo não encontrado: {file_path}")
            print(f"Arquivo não encontrado: {file_path}")
            return ""

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        text = ""

        read_functions = {
            '.txt': read_txt,
            '.docx': read_docx,
            '.pdf': read_pdf,
            '.xlsx': read_excel,
            '.xls': read_excel,
            '.csv': read_csv,
            '.pptx': read_pptx,
            '.html': read_html,
            '.htm': read_html
        }

        if ext in read_functions:
            text = read_functions[ext](file_path)
            if not text.strip():
                logging.warning(f"O arquivo {file_path} está vazio ou não contém texto extraível.")
                print(f"O arquivo {file_path} está vazio ou não contém texto extraível.")
        else:
            logging.error(f"Formato de arquivo não suportado: {ext}")
            print(f"Formato de arquivo não suportado: {ext}")

        return text
    except Exception as e:
        logging.error(f"Erro ao ler o documento: {str(e)}")
        print(f"Erro ao ler o documento: {str(e)}")
        return ""

# Funções de leitura de diferentes formatos de arquivo
def read_txt(file_path):
    """
    Lê arquivos de texto (.txt) com suporte a múltiplas codificações.
    
    Parâmetros:
        file_path (str): Caminho para o arquivo .txt.
    
    Retorna:
        str: Conteúdo do arquivo ou string vazia em caso de erro.
    """
    try:
        encodings = ['utf-8', 'latin-1', 'utf-16']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    logging.info(f"Lendo arquivo TXT com encoding {enc}.")
                    print(f"Lendo arquivo TXT com encoding {enc}.")
                    return f.read()
            except UnicodeDecodeError:
                logging.warning(f"Falha ao ler {file_path} com encoding {enc}. Tentando próximo encoding.")
                print(f"Falha ao ler {file_path} com encoding {enc}. Tentando próximo encoding.")
                continue
            except Exception as e:
                logging.error(f"Erro ao ler arquivo TXT com encoding {enc}: {str(e)}")
                print(f"Erro ao ler arquivo TXT com encoding {enc}: {str(e)}")
                return ""
        logging.error(f"Não foi possível ler o arquivo TXT com as codificações tentadas: {file_path}")
        print(f"Não foi possível ler o arquivo TXT com as codificações tentadas: {file_path}")
        return ""
    except Exception as e:
        logging.error(f"Erro ao ler arquivo TXT: {str(e)}")
        print(f"Erro ao ler arquivo TXT: {str(e)}")
        return ""

def read_docx(file_path):
    """
    Lê arquivos do Word (.docx).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo .docx.
    
    Retorna:
        str: Conteúdo do arquivo ou string vazia em caso de erro.
    """
    try:
        import docx2txt
        logging.info(f"Lendo arquivo DOCX: {file_path}")
        print(f"Lendo arquivo DOCX: {file_path}")
        return docx2txt.process(file_path)
    except Exception as e:
        logging.error(f"Erro ao ler arquivo DOCX: {str(e)}")
        print(f"Erro ao ler arquivo DOCX: {str(e)}")
        return ""

def read_pdf(file_path):
    """
    Lê arquivos PDF (.pdf).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo .pdf.
    
    Retorna:
        str: Conteúdo do arquivo ou string vazia em caso de erro.
    """
    try:
        from pdfminer.high_level import extract_text
        logging.info(f"Lendo arquivo PDF: {file_path}")
        print(f"Lendo arquivo PDF: {file_path}")
        return extract_text(file_path)
    except Exception as e:
        logging.error(f"Erro ao ler arquivo PDF: {str(e)}")
        print(f"Erro ao ler arquivo PDF: {str(e)}")
        return ""

def read_excel(file_path):
    """
    Lê arquivos Excel (.xlsx, .xls).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo Excel.
    
    Retorna:
        str: Conteúdo concatenado das células ou string vazia em caso de erro.
    """
    try:
        import pandas as pd
        logging.info(f"Lendo arquivo Excel: {file_path}")
        print(f"Lendo arquivo Excel: {file_path}")
        df = pd.read_excel(file_path)
        return ' '.join(df.astype(str).stack().tolist())
    except Exception as e:
        logging.error(f"Erro ao ler arquivo Excel: {str(e)}")
        print(f"Erro ao ler arquivo Excel: {str(e)}")
        return ""

def read_csv(file_path):
    """
    Lê arquivos CSV (.csv).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo .csv.
    
    Retorna:
        str: Conteúdo concatenado das células ou string vazia em caso de erro.
    """
    try:
        import pandas as pd
        logging.info(f"Lendo arquivo CSV: {file_path}")
        print(f"Lendo arquivo CSV: {file_path}")
        df = pd.read_csv(file_path)
        return ' '.join(df.astype(str).stack().tolist())
    except Exception as e:
        logging.error(f"Erro ao ler arquivo CSV: {str(e)}")
        print(f"Erro ao ler arquivo CSV: {str(e)}")
        return ""

def read_pptx(file_path):
    """
    Lê arquivos PowerPoint (.pptx).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo .pptx.
    
    Retorna:
        str: Conteúdo das apresentações ou string vazia em caso de erro.
    """
    try:
        from pptx import Presentation
        logging.info(f"Lendo arquivo PPTX: {file_path}")
        print(f"Lendo arquivo PPTX: {file_path}")
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return ' '.join(text_runs)
    except Exception as e:
        logging.error(f"Erro ao ler arquivo PPTX: {str(e)}")
        print(f"Erro ao ler arquivo PPTX: {str(e)}")
        return ""

def read_html(file_path):
    """
    Lê arquivos HTML (.html, .htm).
    
    Parâmetros:
        file_path (str): Caminho para o arquivo HTML.
    
    Retorna:
        str: Texto extraído do HTML ou string vazia em caso de erro.
    """
    try:
        from bs4 import BeautifulSoup
        logging.info(f"Lendo arquivo HTML: {file_path}")
        print(f"Lendo arquivo HTML: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text()
    except Exception as e:
        logging.error(f"Erro ao ler arquivo HTML: {str(e)}")
        print(f"Erro ao ler arquivo HTML: {str(e)}")
        return ""

def main(input_file=None, output_folder=None):
    """
    Função principal que coordena a análise de text mining.
    
    Parâmetros:
        input_file (str): Caminho para o arquivo de entrada.
        output_folder (str): Caminho para a pasta de saída.
    
    Retorna:
        None
    """
    if not input_file or not output_folder:
        root = tk.Tk()
        gui = TextMiningGUI(root)
        root.mainloop()
        input_file = gui.input_file
        output_folder = gui.output_folder

    if not input_file or not output_folder:
        print("Análise cancelada. Nenhum arquivo ou pasta de saída selecionado.")
        return

    setup_logging(output_folder)

    logging.info("Iniciando análise de text mining avançada")
    print("Iniciando análise de text mining avançada")

    try:
        print("Baixando recursos necessários...")
        download_nltk_packages()

        print("Lendo e processando o documento...")
        # Ler e processar o documento
        text = read_document(input_file)
        if not text:
            raise ValueError("O documento está vazio ou não pôde ser lido.")

        from langdetect import detect
        language = detect(text)
        language = 'pt' if language.startswith('pt') else 'en'
        logging.info(f"Idioma detectado: {language}")
        print(f"Idioma detectado: {language}")

        tokens = preprocess_text(text, language)

        # Realizar análises
        print("Realizando análises...")
        word_freq = dict(pd.Series(tokens).value_counts())
        entities = extract_entities(text, language)
        pos_tags = extract_pos_tags(text, language)
        dependencies = dependency_parsing(text, language)
        keywords = keyword_extraction(text, language)
        relationships = extract_relationships(text, language)
        topics = lda_topic_modeling(tokens, language)
        sentiment = sentiment_analysis(text, language)
        stats = text_statistics(text)
        connectives = analyze_connectors(text, language)
        spelling = spelling_correction(text, language)
        readability = readability_scores(text, language)
        dates = extract_dates(text)
        actions = extract_actions_and_responsibles(text, language)
        verb_agreement_errors = check_verb_agreement(text, language)
        person_changes = detect_person_changes(text, language)

        # Gerar visualizações
        print("Gerando visualizações...")
        word_cloud_path = os.path.join(output_folder, 'word_cloud.png')
        generate_word_cloud(word_freq, word_cloud_path)

        entity_network_path = os.path.join(output_folder, 'entity_network.png')
        generate_entity_network(entities, entity_network_path)

        dense_pixel_path = os.path.join(output_folder, 'dense_pixel_display.png')
        generate_dense_pixel_display(text, dense_pixel_path)

        topic_visualization_path = os.path.join(output_folder, 'topic_visualization.png')
        generate_topic_visualization(topics, topic_visualization_path)

        action_flow_path = os.path.join(output_folder, 'action_flow.png')
        generate_action_flow(actions, action_flow_path)

        # Gerar explicações detalhadas para cada seção
        method_explanations = {}
        sections = {
            "Conectores e Preposições Mais Utilizados": "Análise das palavras que conectam ideias e estabelecem relações entre as partes do texto.",
            "Sugestões de Correção Ortográfica": "Identificação de possíveis erros ortográficos e apresentação de sugestões de correção.",
            "Avaliação de Legibilidade": "Utilização de índices de legibilidade para determinar a facilidade de leitura e compreensão do texto.",
            "Frequência de Palavras": "Análise das palavras mais frequentes no texto para identificar os temas centrais.",
            "Nuvem de Palavras": "Visualização gráfica das palavras mais frequentes, onde o tamanho de cada palavra é proporcional à sua frequência no texto.",
            "Entidades Nomeadas": "Processo de identificação e classificação de elementos importantes no texto, como pessoas, organizações e locais.",
            "Rede de Entidades": "Visualização que mostra as relações entre as entidades nomeadas identificadas no texto.",
            "Modelagem de Tópicos": "Utilização de técnicas de modelagem de tópicos para identificar os principais assuntos discutidos no documento.",
            "Visualização de Tópicos": "Representação gráfica dos tópicos identificados e sua relevância no texto.",
            "Análise de Sentimento": "Avaliação da polaridade emocional do texto para determinar se é positivo, negativo ou neutro.",
            "Dense Pixel Display": "Visualização que representa a intensidade e distribuição das palavras ao longo do texto, permitindo identificar padrões de uso e áreas com maior densidade de informação.",
            "Extração de Datas": "Identificação de datas relevantes no documento.",
            "Extração de Ações e Responsáveis": "Mapeamento das ações e seus responsáveis.",
            "Verificação de Concordância Verbal": "Análise da concordância entre sujeito e verbo.",
            "Detecção de Mudanças de Pessoa Gramatical": "Verificação da consistência na pessoa gramatical utilizada.",
            "Fluxo de Ações": "Visualização do encadeamento lógico das ações.",
            "Armazenamento de Dados": "Estruturação das informações em banco de dados.",
            "Part-of-Speech Tagging": "Identificação das classes gramaticais das palavras no texto.",
            "Análise de Dependências": "Análise das relações gramaticais entre as palavras.",
            "Extração de Palavras-Chave": "Identificação das palavras mais relevantes no texto.",
            "Extração de Relações": "Identificação de relações semânticas entre entidades.",
        }
        for section, description in sections.items():
            explanation = generate_method_explanation(section, description)
            method_explanations[section] = explanation

        # Preparar resultados
        analysis_results = {
            'word_frequency': dict(list(word_freq.items())[:20]),  # Top 20 palavras
            'entities': entities[:20],
            'pos_tags': pos_tags[:20],
            'dependencies': dependencies[:20],
            'keywords': keywords,
            'relationships': relationships[:20],
            'topics': topics[:5],
            'sentiment': sentiment,
            'word_cloud_path': word_cloud_path,
            'entity_network_path': entity_network_path,
            'dense_pixel_path': dense_pixel_path,
            'topic_visualization_path': topic_visualization_path,
            'action_flow_path': action_flow_path,
            'text_statistics': stats,
            'connectives': connectives,
            'spelling_corrections': spelling,
            'readability': readability,
            'dates': dates[:10],
            'actions': actions[:10],
            'verb_agreement_errors': verb_agreement_errors,
            'person_changes': person_changes,
            'method_explanations': method_explanations
        }

        # Armazenar dados em banco de dados
        print("Armazenando dados em banco de dados...")
        store_data_in_database(analysis_results, output_folder)

        # Gerar relatório ABNT
        print("Gerando relatório...")
        report_path = os.path.join(output_folder, 'relatorio_analise_abnt.pdf')
        generate_abnt_report(analysis_results, report_path)

        logging.info("Análise concluída com sucesso")
        print("Análise concluída. Os resultados foram salvos na pasta selecionada.")
        logging.info("advanced_text_mining_toolkit_with_LMs - GPT v15.py")
    except Exception as e:
        logging.error(f"Erro durante a análise: {str(e)}")
        print(f"Ocorreu um erro durante a análise: {str(e)}")

def generate_method_explanation(section, description):
    """
    Gera explicações detalhadas para cada método de análise.
    
    Parâmetros:
        section (str): Nome da seção/método.
        description (str): Descrição da seção/método.
    
    Retorna:
        str: Explicação detalhada.
    """
    return f"<b>{section}</b>: {description}\n"

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Advanced Text Mining Analysis Toolkit")
    parser.add_argument('--test', action='store_true', help="Executa os testes internos.")
    parser.add_argument('--input_file', type=str, help="Caminho para o arquivo de entrada.")
    parser.add_argument('--output_folder', type=str, help="Caminho para a pasta de saída.")
    args = parser.parse_args()

    if args.test:
        from tests.test_preprocessing import TestPreprocessing
        from tests.test_analysis import TestAnalysis
        from tests.test_visualization import TestVisualization
        from tests.test_database import TestDatabase
        from tests.test_gui import TestGUI

        setup_logging('test_logs')
        logging.info("Iniciando testes internos.")

        import unittest
        unittest.main(argv=['first-arg-is-ignored'], exit=False)

        print("Todos os testes internos foram executados.")
        logging.info("Todos os testes internos foram executados.")
    else:
        main(args.input_file, args.output_folder)
